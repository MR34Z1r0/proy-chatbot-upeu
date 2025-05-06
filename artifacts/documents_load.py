import os
import logging
import boto3
import hashlib
import requests
import unicodedata
from uuid import uuid4
from dotenv import load_dotenv

# Procesamiento de documentos
from pptx import Presentation as PptxPresentation
from docx import Document as DocxDocument
from openpyxl import load_workbook

# Langchain y Pinecone
from langchain.document_loaders import PyPDFLoader
from langchain.text_splitter import CharacterTextSplitter, TokenTextSplitter
from langchain.embeddings.bedrock import BedrockEmbeddings
from langchain_core.documents import Document
from langchain_pinecone import PineconeVectorStore
from pinecone import Pinecone

# Utilidades personalizadas
from artifacts.dynamodb_utils import DYNAMODB_RESOURCES, DYNAMODB_LIBRARY
from artifacts.pinecone_utils import delete_vectors_by_file_hash

# Cargar variables de entorno
load_dotenv()

# Configuración de logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

# -------------------------------
# Configuración de AWS y recursos
# -------------------------------

# Establecer perfil de sesión para boto3
boto3.setup_default_session(profile_name='prd-sofia-admin')

# Clientes de AWS
s3 = boto3.client('s3')

# Claves de AWS desde .env
AWS_ACCESS_KEY_ID = os.getenv("AWS_ACCESS_KEY_ID")
AWS_SECRET_ACCESS_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
# Inicializar recursos de DynamoDB
dynamodb_resources = DYNAMODB_RESOURCES(AWS_ACCESS_KEY_ID, AWS_SECRET_ACCESS_KEY)
dynamodb_library = DYNAMODB_LIBRARY(AWS_ACCESS_KEY_ID, AWS_SECRET_ACCESS_KEY)

# -------------------------------
# Configuración de Pinecone
# -------------------------------
logging.info("Inicializando conexión con Pinecone.")

PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_INDEX_NAME = os.getenv("PINECONE_INDEX_NAME")
EMBEDDINGS_MODEL_ID = os.getenv("EMBEDDINGS_MODEL_ID")
DIMENSION = 1024

pc = Pinecone(api_key=PINECONE_API_KEY)
index = pc.Index(PINECONE_INDEX_NAME)

# -------------------------------
# Parámetros generales del sistema
# -------------------------------
bucket_name = 'datalake-cls-509399624591-landing-s3-bucket'
files_table_name = 'db_learning_resources'
hash_table_name = 'db_learning_resources_hash'

# Carpeta local para descargas
download_folder = "./downloads"
os.makedirs(download_folder, exist_ok=True)

# Funciones auxiliares
def sanitize_filename(filename):
    """Limpia caracteres especiales y espacios en el nombre del archivo."""
    filename = ''.join(c for c in unicodedata.normalize('NFD', filename) if unicodedata.category(c) != 'Mn')
    return filename.replace(' ', '_')

def download_file_from_gdrive(file_name, gdrive_id):
    """Descarga un archivo desde Google Drive y lo guarda localmente."""
    url = f"https://drive.google.com/uc?export=download&id={gdrive_id}"
    file_path = os.path.join(download_folder, file_name)
    logging.info(f"Descargando {file_name} desde Google Drive.")
    response = requests.get(url, stream=True)
    response.raise_for_status()
    with open(file_path, 'wb') as f:
        for chunk in response.iter_content(chunk_size=8192):
            f.write(chunk)
    return file_path

def generate_file_hash(file_path):
    """Genera un hash SHA256 para el archivo dado."""
    logging.info("Generando hash para el archivo en memoria")
    sha256_hash = hashlib.sha256()
    with open(file_path, 'rb') as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()

def upload_to_s3(file_path, bucket, object_name):
    """Sube un archivo local al bucket de S3."""
    full_object_name = f"SOFIA_FILE/PLANIFICACION/AV_Recursos/{object_name}"
    logging.info(f"Subiendo archivo a S3 bucket {bucket} en la ruta {full_object_name}.")
    s3.upload_file(file_path, bucket, full_object_name)
    logging.info(f"Archivo {full_object_name} subido exitosamente.")
    return f"s3://{bucket}/{full_object_name}"

def delete_from_s3(bucket, object_name):
    """Elimina un archivo de un bucket de S3."""
    full_object_name = f"SOFIA_FILE/PLANIFICACION/AV_Recursos/{object_name}"
    logging.info(f"Eliminando archivo de S3 bucket {bucket} en la ruta {full_object_name}.")
    s3.delete_object(Bucket=bucket, Key=full_object_name)
    logging.info(f"Archivo {full_object_name} eliminado exitosamente.")
    return f"Archivo s3://{bucket}/{full_object_name} eliminado."

# Funciones para los tipos de archivo
def process_file_pdf(metadata, file_path):
    """Procesa un archivo PDF."""
    loader = PyPDFLoader(file_path=file_path)
    data = loader.load()
    text_splitter = CharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=400,
        chunk_overlap=20
    )
    docs = text_splitter.split_documents(documents=data)
    return add_to_pinecone(metadata, docs)

def process_file_pptx(metadata, file_path):
    """Procesa un archivo PPTX."""
    prs = PptxPresentation(file_path)
    all_slides = []
    for slide in prs.slides:
        list_slide = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                list_slide.append(shape.text)
        all_slides.append("\n".join(list_slide))

    documents = [Document(page_content=slide, metadata=metadata) for slide in all_slides]
    text_splitter = CharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=400,
        chunk_overlap=20
    )
    docs = text_splitter.split_documents(documents=documents)
    return add_to_pinecone(metadata, docs)

def process_file_docx(metadata, file_path):
    """Procesa un archivo DOCX."""
    doc = DocxDocument(file_path)
    all_page = [para.text for para in doc.paragraphs]

    documents = [Document(page_content=content, metadata=metadata) for content in all_page]
    text_splitter = CharacterTextSplitter.from_tiktoken_encoder(
        chunk_size=400,
        chunk_overlap=20
    )
    docs = text_splitter.split_documents(documents=documents)
    return add_to_pinecone(metadata, docs)

def process_file_xlsx(metadata, file_path):
    """Procesa un archivo XLSX dividiendo el contenido estrictamente para cumplir con el límite de tokens."""
    try:
        workbook = load_workbook(file_path)
        all_sheets_data = []
        max_tokens = 8000  # Límite seguro de tokens (menor que 8192)

        # Itera sobre todas las hojas del archivo
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_data = []

            # Lee todas las filas y columnas de la hoja
            for row in sheet.iter_rows(values_only=True):
                sheet_data.append(row)

            # Convierte los datos de la hoja a un string
            content = f"Sheet Name: {sheet_name}\nData:\n{sheet_data}"

            # Divide el contenido estrictamente en fragmentos de 8000 tokens
            text_splitter = TokenTextSplitter(chunk_size=max_tokens, chunk_overlap=0)
            chunks = text_splitter.split_text(content)

            for chunk in chunks:
                all_sheets_data.append(Document(page_content=chunk, metadata=metadata))

        # Agrega los documentos a Pinecone
        return add_to_pinecone(metadata, all_sheets_data)

    except Exception as e:
        logging.error(f"Error procesando archivo XLSX: {e}")
        return []

def add_to_pinecone(metadata, docs):
    """Agrega documentos a Pinecone y devuelve los IDs generados."""
    embeddings = BedrockEmbeddings(
        model_id=EMBEDDINGS_MODEL_ID,
        credentials_profile_name="prd-sofia-admin",
        region_name="us-west-2"
    )
    vector_store = PineconeVectorStore(index=index, embedding=embeddings)
    uuids = [str(uuid4()) for _ in range(len(docs))] 
    updated_docs = []
    for doc, uuid in zip(docs, uuids):
        doc.metadata.update(metadata)
        updated_docs.append(doc)

    vector_store.add_documents(documents=updated_docs, ids=uuids)
    logging.info("Documentos agregados a Pinecone.")
    return uuids

# Añadir archivo
def add_file(event):
    """Procesa el evento para descargar, subir a S3 y registrar en DynamoDB."""
    try:
        logging.info("Procesando evento.")

        # Descarga el archivo desde Google Drive
        file_name = sanitize_filename(event['TituloRecurso'])
        gdrive_id = event['DriveId']
        file_path = download_file_from_gdrive(file_name, gdrive_id)

        # Genera el hash del archivo descargado
        file_hash = generate_file_hash(file_path)

        # Sube el archivo a S3
        s3_path = upload_to_s3(file_path, bucket_name, file_name)

        # Construye el objeto para DynamoDB
        dynamodb_resource_obj = {
            'resource_id': event['RecursoDidacticoId'],
            'resource_title': event['TituloRecurso'],
            'drive_id': gdrive_id,
            'file_hash': file_hash,
            's3_path': s3_path,
            #'pinecone_ids': []  # Lista para almacenar los IDs de Pinecone
        }
        
        print(dynamodb_resource_obj)
        # Sube los datos a DynamoDB (si ya existe el file_hash, ya no lo sube a Pinecone porque ya fue subido anteriormente)
        result = dynamodb_resources.upload_in_resources(dynamodb_resource_obj)
        if result:
            logging.info("Procesando el archivo para indexación en Pinecone.")
            pinecone_ids = []

            if file_name.lower().endswith('pdf'):
                pinecone_ids = process_file_pdf(dynamodb_resource_obj, file_path)
            elif file_name.lower().endswith('pptx'):
                pinecone_ids = process_file_pptx(dynamodb_resource_obj, file_path)
            elif file_name.lower().endswith('docx'):
                pinecone_ids = process_file_docx(dynamodb_resource_obj, file_path)
            elif file_name.lower().endswith('xlsx'):
                pinecone_ids = process_file_xlsx(dynamodb_resource_obj, file_path)
            else: 
                logging.warning(f"Formato de archivo no compatible: {file_name}")

            # Actualiza los IDs de Pinecone en DynamoDB
            success = dynamodb_resources.update_in_resources_to_pinecone_ids(event['RecursoDidacticoId'], pinecone_ids)
            if success:
                print("Pinecone IDs actualizados correctamente.")
            else:
                print("Error al actualizar los Pinecone IDs.")

            logging.info("Evento procesado exitosamente.")

        # Eliminar el archivo descargado
        os.remove(file_path)
        logging.info(f"Archivo temporal {file_path} eliminado.")

    except Exception as e:
        logging.error(f"Error procesando el evento: {e}")

def delete_file(event):
    resource_id = event['RecursoDidacticoId']
    file_name = sanitize_filename(event['TituloRecurso'])
    silabus_id = event['SilaboEventoId']
    
    # Eliminar de DynamoDB (resources y resources_hash) y Pinecone
    dynamodb_resources.delete_resource_and_vectors(resource_id, index)

    # Eliminar del listado de resources en el silabus
    dynamodb_library.remove_resource_from_library(silabus_id, resource_id)

    # Eliminar de S3
    delete_from_s3(bucket_name, file_name)